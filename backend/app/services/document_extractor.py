from pathlib import Path

from app.config import settings


def _resolve_path(file_path: str) -> str:
    """Convert public /uploads/... paths to local filesystem paths."""
    if file_path.startswith("/uploads/"):
        return str(Path(settings.upload_dir) / file_path[len("/uploads/"):])
    return file_path


def extract_text(file_path: str) -> str:
    file_path = _resolve_path(file_path)
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".docx":
        return _extract_docx(file_path)
    if ext in (".ppt", ".pptx"):
        return _extract_pptx(file_path)
    if ext in (".xls", ".xlsx"):
        return _extract_xlsx(file_path)
    if ext in (".txt", ".md", ".py", ".json", ".csv"):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    return ""


def _extract_pdf(path: str) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:
        return f"[PDF extraction error: {exc}]"


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as exc:
        return f"[DOCX extraction error: {exc}]"


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
        return "\n".join(lines)
    except Exception as exc:
        return f"[PPTX extraction error: {exc}]"


def _extract_xlsx(path: str) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                lines.append(" ".join(str(cell) for cell in row if cell is not None))
        return "\n".join(lines)
    except Exception as exc:
        return f"[XLSX extraction error: {exc}]"
